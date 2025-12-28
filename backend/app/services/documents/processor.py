"""
OfflineRAG - Document Processing Service
=========================================

Handles extraction, parsing, and chunking of various document types.
Supports: Text, PDF, Word, Excel, CSV, Images, Audio, Video
"""

import os
import asyncio
import hashlib
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple, AsyncGenerator
from datetime import datetime
import chardet
import mimetypes

from loguru import logger

from app.core.config import settings
from app.models.schemas import (
    Document, DocumentChunk, DocumentMetadata, 
    DocumentType, DocumentStatus
)


class DocumentProcessor:
    """Main document processing orchestrator."""
    
    def __init__(self):
        self._extractors: Dict[DocumentType, 'BaseExtractor'] = {}
        self._chunker = TextChunker()
        self._initialized = False
        self._cancel_flags: Dict[str, bool] = {}
    
    async def initialize(self):
        """Initialize all extractors."""
        if self._initialized:
            return
        
        logger.info("Initializing document processors...")
        
        # Register extractors
        self._extractors = {
            DocumentType.TEXT: TextExtractor(),
            DocumentType.PDF: PDFExtractor(),
            DocumentType.WORD: WordExtractor(),
            DocumentType.EXCEL: ExcelExtractor(),
            DocumentType.CSV: CSVExtractor(),
            DocumentType.IMAGE: ImageExtractor(),
            DocumentType.AUDIO: AudioExtractor(),
            DocumentType.VIDEO: VideoExtractor(),
            DocumentType.POWERPOINT: PowerPointExtractor(),
        }
        
        # Initialize each extractor
        for doc_type, extractor in self._extractors.items():
            try:
                await extractor.initialize()
                logger.debug(f"Initialized {doc_type.value} extractor")
            except Exception as e:
                logger.warning(f"Failed to initialize {doc_type.value} extractor: {e}")
        
        self._initialized = True
        logger.info("Document processors initialized")
    
    def detect_type(self, filename: str) -> DocumentType:
        """Detect document type from filename extension."""
        ext = Path(filename).suffix.lower()
        
        if ext in settings.SUPPORTED_TEXT_EXTENSIONS:
            return DocumentType.TEXT
        elif ext in [".pdf"]:
            return DocumentType.PDF
        elif ext in [".docx", ".doc"]:
            return DocumentType.WORD
        elif ext in [".xlsx", ".xls"]:
            return DocumentType.EXCEL
        elif ext in [".csv"]:
            return DocumentType.CSV
        elif ext in settings.SUPPORTED_IMAGE_EXTENSIONS:
            return DocumentType.IMAGE
        elif ext in settings.SUPPORTED_AUDIO_EXTENSIONS:
            return DocumentType.AUDIO
        elif ext in settings.SUPPORTED_VIDEO_EXTENSIONS:
            return DocumentType.VIDEO
        elif ext in [".pptx", ".ppt"]:
            return DocumentType.POWERPOINT
        else:
            return DocumentType.UNKNOWN
    
    def request_cancel(self, document_id: str):
        """Request cancellation of document processing."""
        self._cancel_flags[document_id] = True
        logger.info(f"Cancellation requested for document {document_id}")
    
    def is_cancelled(self, document_id: str) -> bool:
        """Check if processing is cancelled."""
        return self._cancel_flags.get(document_id, False)
    
    def clear_cancel_flag(self, document_id: str):
        """Clear cancellation flag."""
        self._cancel_flags.pop(document_id, None)
    
    async def process_file(
        self,
        file_path: Path,
        filename: str,
        document_id: str,
        progress_callback: Optional[callable] = None
    ) -> Tuple[str, List[DocumentChunk], DocumentMetadata]:
        """
        Process a file and return extracted content, chunks, and metadata.
        
        Args:
            file_path: Path to the file
            filename: Original filename
            document_id: Unique document ID
            progress_callback: Optional callback for progress updates
            
        Returns:
            Tuple of (raw_text, chunks, metadata)
        """
        doc_type = self.detect_type(filename)
        extractor = self._extractors.get(doc_type)
        
        if not extractor:
            raise ValueError(f"No extractor available for type: {doc_type}")
        
        # Check for cancellation
        if self.is_cancelled(document_id):
            self.clear_cancel_flag(document_id)
            raise asyncio.CancelledError("Processing cancelled by user")
        
        # Extract content
        if progress_callback:
            await progress_callback(0.1, "Extracting content...")
        
        raw_text, extract_metadata = await extractor.extract(file_path)
        
        if self.is_cancelled(document_id):
            self.clear_cancel_flag(document_id)
            raise asyncio.CancelledError("Processing cancelled by user")
        
        # Create metadata
        file_stat = file_path.stat()
        metadata = DocumentMetadata(
            filename=filename,
            file_type=doc_type,
            file_size=file_stat.st_size,
            word_count=len(raw_text.split()) if raw_text else 0,
            **extract_metadata
        )
        
        if progress_callback:
            await progress_callback(0.5, "Chunking content...")
        
        # Chunk the content
        chunks = await self._chunker.chunk_text(
            raw_text,
            document_id=document_id,
            metadata={
                "filename": filename,
                "file_type": doc_type.value,
            }
        )
        
        if self.is_cancelled(document_id):
            self.clear_cancel_flag(document_id)
            raise asyncio.CancelledError("Processing cancelled by user")
        
        if progress_callback:
            await progress_callback(0.9, "Finalizing...")
        
        self.clear_cancel_flag(document_id)
        return raw_text, chunks, metadata


class BaseExtractor:
    """Base class for document extractors."""
    
    async def initialize(self):
        """Initialize the extractor."""
        pass
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text content from file.
        
        Returns:
            Tuple of (text_content, metadata_dict)
        """
        raise NotImplementedError


class TextExtractor(BaseExtractor):
    """Extractor for plain text files."""
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        # Detect encoding
        with open(file_path, 'rb') as f:
            raw = f.read()
            result = chardet.detect(raw)
            encoding = result['encoding'] or 'utf-8'
        
        # Read content
        content = raw.decode(encoding, errors='replace')
        
        metadata = {
            "encoding": encoding,
            "line_count": content.count('\n') + 1,
        }
        
        return content, metadata


class PDFExtractor(BaseExtractor):
    """Extractor for PDF documents."""
    
    def __init__(self):
        self._use_ocr = settings.OCR_ENABLED
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        import fitz  # PyMuPDF
        
        texts = []
        page_count = 0
        
        # Open PDF
        doc = fitz.open(str(file_path))
        page_count = len(doc)
        
        for page_num, page in enumerate(doc):
            # Try text extraction first
            text = page.get_text()
            
            # If no text and OCR enabled, use OCR
            if not text.strip() and self._use_ocr:
                text = await self._ocr_page(page)
            
            if text:
                texts.append(f"[Page {page_num + 1}]\n{text}")
        
        doc.close()
        
        metadata = {
            "page_count": page_count,
        }
        
        return "\n\n".join(texts), metadata
    
    async def _ocr_page(self, page) -> str:
        """OCR a PDF page."""
        try:
            import pytesseract
            from PIL import Image
            import io
            
            # Render page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))
            
            # OCR
            text = pytesseract.image_to_string(image, lang=settings.OCR_LANGUAGE)
            return text
        except Exception as e:
            logger.warning(f"OCR failed: {e}")
            return ""


class WordExtractor(BaseExtractor):
    """Extractor for Word documents."""
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        from docx import Document as DocxDocument
        
        doc = DocxDocument(str(file_path))
        
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(" | ".join(row_text))
            texts.append("\n".join(table_text))
        
        metadata = {
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
        }
        
        return "\n\n".join(texts), metadata


class ExcelExtractor(BaseExtractor):
    """
    Enhanced Excel extractor with schema-aware parsing.
    
    Features:
    - Header detection
    - Data type inference
    - Column summaries
    - Statistical insights
    - Efficient chunking for large files
    """
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        import pandas as pd
        
        texts = []
        total_rows = 0
        total_cols = 0
        all_columns = []
        
        # Read all sheets
        xlsx = pd.ExcelFile(file_path)
        sheet_count = len(xlsx.sheet_names)
        
        for sheet_name in xlsx.sheet_names:
            df = pd.read_excel(xlsx, sheet_name=sheet_name)
            total_rows += len(df)
            total_cols = max(total_cols, len(df.columns))
            all_columns.extend(df.columns.tolist())
            
            sheet_text = await self._process_dataframe(df, f"Sheet: {sheet_name}")
            texts.append(sheet_text)
        
        metadata = {
            "sheet_count": sheet_count,
            "total_rows": total_rows,
            "total_columns": total_cols,
            "all_columns": list(set(all_columns)),
        }
        
        return "\n\n".join(texts), metadata
    
    async def _process_dataframe(self, df, title: str) -> str:
        """Process a DataFrame with schema awareness."""
        import pandas as pd
        
        parts = [f"[{title}]"]
        
        # Schema information
        schema_info = []
        for col in df.columns:
            dtype = str(df[col].dtype)
            null_count = df[col].isnull().sum()
            unique_count = df[col].nunique()
            schema_info.append(f"  - {col}: {dtype} (nulls: {null_count}, unique: {unique_count})")
        
        parts.append("Schema:\n" + "\n".join(schema_info))
        
        # Statistical summary for numeric columns
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            stats = df[numeric_cols].describe().to_string()
            parts.append(f"Statistics:\n{stats}")
        
        # Sample data (first and last rows)
        if len(df) <= 20:
            parts.append(f"Data:\n{df.to_string(index=False)}")
        else:
            head = df.head(10).to_string(index=False)
            tail = df.tail(5).to_string(index=False)
            parts.append(f"First 10 rows:\n{head}")
            parts.append(f"Last 5 rows:\n{tail}")
            parts.append(f"(Total: {len(df)} rows)")
        
        return "\n\n".join(parts)


class CSVExtractor(BaseExtractor):
    """
    Enhanced CSV extractor with schema-aware parsing.
    
    Features:
    - Automatic delimiter detection
    - Header detection
    - Data type inference
    - Column summaries
    - Row-level semantic chunks
    - Analytical query support
    """
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        import pandas as pd
        import csv
        
        # Detect delimiter
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            sample = f.read(4096)
            try:
                dialect = csv.Sniffer().sniff(sample)
                delimiter = dialect.delimiter
            except:
                delimiter = ','
        
        # Read CSV
        df = pd.read_csv(file_path, delimiter=delimiter)
        
        parts = ["[CSV Data Analysis]"]
        
        # Schema information
        schema_info = []
        for col in df.columns:
            dtype = str(df[col].dtype)
            null_count = df[col].isnull().sum()
            unique_count = df[col].nunique()
            sample_vals = df[col].dropna().head(3).tolist()
            sample_str = str(sample_vals)[:50]
            schema_info.append(f"  - {col}: {dtype} (nulls: {null_count}, unique: {unique_count}, sample: {sample_str})")
        
        parts.append("Columns:\n" + "\n".join(schema_info))
        
        # Numeric statistics
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            stats = df[numeric_cols].describe()
            parts.append(f"Numeric Statistics:\n{stats.to_string()}")
        
        # Categorical summaries
        cat_cols = df.select_dtypes(include=['object', 'category']).columns
        for col in cat_cols[:5]:  # Limit to first 5 categorical columns
            value_counts = df[col].value_counts().head(10)
            parts.append(f"Top values for '{col}':\n{value_counts.to_string()}")
        
        # Data preview
        if len(df) <= 50:
            parts.append(f"All Data:\n{df.to_string(index=False)}")
        else:
            parts.append(f"Sample (first 20 rows):\n{df.head(20).to_string(index=False)}")
            parts.append(f"(Total: {len(df)} rows)")
        
        metadata = {
            "row_count": len(df),
            "column_count": len(df.columns),
            "columns": list(df.columns),
            "delimiter": delimiter,
            "numeric_columns": list(numeric_cols),
            "categorical_columns": list(cat_cols),
            "has_nulls": df.isnull().any().any(),
        }
        
        return "\n\n".join(parts), metadata


class ImageExtractor(BaseExtractor):
    """Extractor for images using OCR and BLIP captioning."""
    
    _tesseract_available = False
    _blip_model = None
    _blip_processor = None
    _blip_available = False
    
    async def initialize(self):
        # Initialize Tesseract OCR
        if settings.OCR_ENABLED:
            try:
                import pytesseract
                
                # Try to set Tesseract command
                tesseract_cmd = settings.get_tesseract_cmd()
                if tesseract_cmd:
                    pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
                    logger.info(f"Using Tesseract at: {tesseract_cmd}")
                
                # Test if tesseract is actually available
                version = pytesseract.get_tesseract_version()
                ImageExtractor._tesseract_available = True
                logger.info(f"Tesseract OCR v{version} is available")
            except ImportError:
                logger.warning("pytesseract not installed - OCR disabled. Install with: pip install pytesseract")
            except Exception as e:
                logger.warning(f"Tesseract not available: {e}. Install Tesseract-OCR from https://github.com/UB-Mannheim/tesseract/wiki")
        
        # Initialize BLIP for image captioning
        try:
            from transformers import BlipProcessor, BlipForConditionalGeneration
            import torch
            
            model_name = "Salesforce/blip-image-captioning-base"
            logger.info(f"Loading BLIP model: {model_name}")
            
            ImageExtractor._blip_processor = await asyncio.to_thread(
                BlipProcessor.from_pretrained, model_name, use_fast=True
            )
            ImageExtractor._blip_model = await asyncio.to_thread(
                BlipForConditionalGeneration.from_pretrained, model_name
            )
            
            # Move to GPU if available
            device = "cuda" if torch.cuda.is_available() else "cpu"
            ImageExtractor._blip_model = ImageExtractor._blip_model.to(device)
            ImageExtractor._blip_available = True
            logger.info(f"BLIP image captioning loaded on {device}")
        except Exception as e:
            logger.warning(f"BLIP captioning not available: {e}")
    
    async def _generate_caption(self, image) -> str:
        """Generate a caption for the image using BLIP."""
        if not ImageExtractor._blip_available:
            return ""
        
        try:
            import torch
            
            device = next(ImageExtractor._blip_model.parameters()).device
            
            # Prepare image
            inputs = ImageExtractor._blip_processor(image, return_tensors="pt").to(device)
            
            # Generate caption
            with torch.no_grad():
                out = await asyncio.to_thread(
                    ImageExtractor._blip_model.generate,
                    **inputs,
                    max_new_tokens=100
                )
            
            caption = ImageExtractor._blip_processor.decode(out[0], skip_special_tokens=True)
            return caption
        except Exception as e:
            logger.warning(f"Failed to generate image caption: {e}")
            return ""
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        from PIL import Image
        
        image = Image.open(file_path)
        
        # Get image info
        width, height = image.size
        
        # Convert to RGB for processing
        if image.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'P':
                image = image.convert('RGBA')
            background.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
            rgb_image = background
        elif image.mode != 'RGB':
            rgb_image = image.convert('RGB')
        else:
            rgb_image = image
        
        # Perform OCR if available
        text = ""
        if ImageExtractor._tesseract_available:
            try:
                import pytesseract
                text = await asyncio.to_thread(
                    pytesseract.image_to_string,
                    rgb_image,
                    lang=settings.OCR_LANGUAGE
                )
                text = text.strip()
            except Exception as e:
                logger.warning(f"OCR failed for image: {e}")
                text = ""
        
        # Generate BLIP caption (for images without text or as additional context)
        caption = ""
        if ImageExtractor._blip_available:
            caption = await self._generate_caption(rgb_image)
        
        # Image description
        description = f"[Image: {width}x{height} pixels, format: {image.format or file_path.suffix.upper()}]"
        
        metadata = {
            "width": width,
            "height": height,
            "format": image.format or file_path.suffix.replace('.', '').upper(),
            "mode": image.mode,
            "has_text": bool(text),
            "caption": caption,
        }
        
        # Build full text output
        parts = [description]
        
        if caption:
            parts.append(f"\nImage description: {caption}")
        
        if text:
            parts.append(f"\nExtracted text:\n{text}")
        elif not caption:
            parts.append("\n[No text detected in image]")
        
        full_text = "\n".join(parts)
        
        return full_text, metadata


class AudioExtractor(BaseExtractor):
    """Extractor for audio files using Whisper."""
    
    def __init__(self):
        self._model = None
    
    async def initialize(self):
        try:
            import whisper
            self._model = whisper.load_model(settings.WHISPER_MODEL)
            logger.info(f"Loaded Whisper model: {settings.WHISPER_MODEL}")
        except Exception as e:
            logger.warning(f"Failed to load Whisper model: {e}")
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        if not self._model:
            return "[Audio transcription not available]", {}
        
        import whisper
        
        # Transcribe
        result = await asyncio.to_thread(
            self._model.transcribe,
            str(file_path),
            language=settings.WHISPER_LANGUAGE
        )
        
        text = result["text"]
        
        # Get duration
        from pydub import AudioSegment
        audio = AudioSegment.from_file(str(file_path))
        duration = len(audio) / 1000  # seconds
        
        metadata = {
            "duration": duration,
            "language": result.get("language", "unknown"),
        }
        
        return text, metadata


class VideoExtractor(BaseExtractor):
    """
    Enhanced video extractor with multimodal processing.
    
    Features:
    - Audio track extraction and transcription
    - Intelligent frame sampling (time-based)
    - OCR on key frames
    - Semantic frame descriptions
    - Partial success handling
    """
    
    def __init__(self):
        self._whisper_model = None
        self._frame_interval = 5  # Sample every 5 seconds
        self._max_frames = 20  # Maximum frames to process
    
    async def initialize(self):
        try:
            import whisper
            self._whisper_model = whisper.load_model(settings.WHISPER_MODEL)
            logger.info("Video extractor: Whisper model loaded")
        except Exception as e:
            logger.warning(f"Failed to load Whisper for video: {e}")
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        from moviepy.editor import VideoFileClip
        import tempfile
        
        texts = []
        metadata = {}
        audio_success = False
        frames_success = False
        
        clip = None
        try:
            clip = VideoFileClip(str(file_path))
            duration = clip.duration
            metadata['duration'] = duration
            metadata['has_audio'] = clip.audio is not None
            metadata['size'] = (clip.w, clip.h) if hasattr(clip, 'w') else None
            
            # Extract and transcribe audio
            if clip.audio and self._whisper_model:
                try:
                    audio_text = await self._extract_audio_transcript(clip)
                    if audio_text:
                        texts.append(f"[Audio Transcription]\n{audio_text}")
                        audio_success = True
                except Exception as e:
                    logger.warning(f"Audio extraction failed: {e}")
            
            # Sample and process frames
            try:
                frame_texts = await self._process_frames(clip, duration)
                if frame_texts:
                    texts.append("[Visual Analysis]\n" + "\n".join(frame_texts))
                    frames_success = True
            except Exception as e:
                logger.warning(f"Frame processing failed: {e}")
            
            metadata['audio_extracted'] = audio_success
            metadata['frames_extracted'] = frames_success
            
        except Exception as e:
            logger.error(f"Video processing error: {e}")
            texts.append(f"[Video: Processing partially failed - {str(e)[:50]}]")
        finally:
            if clip:
                try:
                    clip.close()
                except:
                    pass
        
        if not texts:
            texts.append(f"[Video: {metadata.get('duration', 0):.1f} seconds]")
        
        return "\n\n".join(texts), metadata
    
    async def _extract_audio_transcript(self, clip) -> Optional[str]:
        """Extract and transcribe audio from video."""
        import tempfile
        
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            try:
                await asyncio.to_thread(
                    clip.audio.write_audiofile,
                    tmp.name,
                    verbose=False,
                    logger=None
                )
                
                result = await asyncio.to_thread(
                    self._whisper_model.transcribe,
                    tmp.name,
                    language=settings.WHISPER_LANGUAGE
                )
                return result.get('text', '')
            finally:
                try:
                    os.unlink(tmp.name)
                except:
                    pass
    
    async def _process_frames(self, clip, duration: float) -> List[str]:
        """Sample and process key frames from video with OCR and BLIP captioning."""
        from PIL import Image
        import numpy as np
        
        frame_texts = []
        
        # Calculate frame timestamps
        num_frames = min(self._max_frames, int(duration / self._frame_interval) + 1)
        timestamps = [i * self._frame_interval for i in range(num_frames)]
        timestamps = [t for t in timestamps if t < duration]
        
        for i, timestamp in enumerate(timestamps):
            try:
                # Extract frame
                frame = await asyncio.to_thread(clip.get_frame, timestamp)
                image = Image.fromarray(np.uint8(frame))
                
                # Try OCR on frame
                ocr_text = await self._ocr_frame(image)
                
                # Try BLIP captioning
                caption = await self._caption_frame(image)
                
                # Generate frame description
                time_str = f"{int(timestamp//60)}:{int(timestamp%60):02d}"
                
                parts = [f"[{time_str}]"]
                
                if caption:
                    parts.append(f"Scene: {caption}")
                
                if ocr_text and len(ocr_text.strip()) > 10:
                    parts.append(f"Text: {ocr_text.strip()[:200]}")
                
                if len(parts) > 1:
                    frame_texts.append(" | ".join(parts))
                else:
                    w, h = image.size
                    frame_texts.append(f"[{time_str}] Frame {i+1}: {w}x{h}")
                    
            except Exception as e:
                logger.debug(f"Frame {i} processing failed: {e}")
                continue
        
        return frame_texts
    
    async def _caption_frame(self, image) -> str:
        """Generate BLIP caption for a video frame."""
        if not ImageExtractor._blip_available:
            return ""
        
        try:
            import torch
            
            device = next(ImageExtractor._blip_model.parameters()).device
            inputs = ImageExtractor._blip_processor(image, return_tensors="pt").to(device)
            
            with torch.no_grad():
                out = await asyncio.to_thread(
                    ImageExtractor._blip_model.generate,
                    **inputs,
                    max_new_tokens=50
                )
            
            caption = ImageExtractor._blip_processor.decode(out[0], skip_special_tokens=True)
            return caption
        except Exception as e:
            logger.debug(f"Frame captioning failed: {e}")
            return ""
    
    async def _ocr_frame(self, image) -> str:
        """Perform OCR on a video frame."""
        try:
            import pytesseract
            text = await asyncio.to_thread(
                pytesseract.image_to_string,
                image,
                lang=settings.OCR_LANGUAGE
            )
            return text
        except Exception:
            return ""


class PowerPointExtractor(BaseExtractor):
    """Extractor for PowerPoint presentations."""
    
    async def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        from pptx import Presentation
        
        prs = Presentation(str(file_path))
        
        texts = []
        slide_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_count += 1
            slide_texts = [f"[Slide {slide_num}]"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            texts.append("\n".join(slide_texts))
        
        metadata = {
            "slide_count": slide_count,
        }
        
        return "\n\n".join(texts), metadata


class TextChunker:
    """Intelligent text chunking with semantic boundary preservation."""
    
    def __init__(
        self,
        chunk_size: int = None,
        chunk_overlap: int = None,
        min_chunk_size: int = None
    ):
        self.chunk_size = chunk_size or settings.CHUNK_SIZE
        self.chunk_overlap = chunk_overlap or settings.CHUNK_OVERLAP
        self.min_chunk_size = min_chunk_size or settings.MIN_CHUNK_SIZE
    
    async def chunk_text(
        self,
        text: str,
        document_id: str,
        metadata: Dict[str, Any] = None
    ) -> List[DocumentChunk]:
        """
        Split text into overlapping chunks while preserving semantic boundaries.
        """
        if not text or not text.strip():
            return []
        
        chunks = []
        
        # Split into paragraphs first
        paragraphs = text.split('\n\n')
        
        current_chunk = ""
        chunk_index = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # If adding this paragraph exceeds chunk size
            if len(current_chunk) + len(para) > self.chunk_size:
                # Save current chunk if it meets minimum size
                if len(current_chunk) >= self.min_chunk_size:
                    chunks.append(DocumentChunk(
                        document_id=document_id,
                        content=current_chunk.strip(),
                        chunk_index=chunk_index,
                        metadata={
                            **(metadata or {}),
                            "char_count": len(current_chunk),
                        }
                    ))
                    chunk_index += 1
                    
                    # Start new chunk with overlap
                    overlap_text = current_chunk[-self.chunk_overlap:] if len(current_chunk) > self.chunk_overlap else ""
                    current_chunk = overlap_text + " " + para
                else:
                    current_chunk += " " + para
            else:
                current_chunk = (current_chunk + " " + para).strip()
        
        # Don't forget the last chunk
        if current_chunk.strip() and len(current_chunk) >= self.min_chunk_size:
            chunks.append(DocumentChunk(
                document_id=document_id,
                content=current_chunk.strip(),
                chunk_index=chunk_index,
                metadata={
                    **(metadata or {}),
                    "char_count": len(current_chunk),
                }
            ))
        
        return chunks


# Singleton instance
document_processor = DocumentProcessor()
